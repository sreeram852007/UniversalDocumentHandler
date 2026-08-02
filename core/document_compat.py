"""
Document Compatibility Module - Handles various document formats with fallback methods
"""

import zipfile
import xml.etree.ElementTree as ET
from PyQt6.QtCore import Qt


class DocumentCompatibility:
    """Superior document handling with fallback methods"""
    
    @staticmethod
    def render_docx_with_formatting(file_path):
        """
        Render DOCX with HTML formatting for better display
        Preserves: Headings, Bold, Italic, Underline, Lists, Tables
        """
        try:
            from docx import Document
            doc = Document(file_path)
            
            html_parts = []
            
            # Process paragraphs
            for para in doc.paragraphs:
                # Get paragraph style
                style_name = para.style.name if para.style else ""
                para_html = DocumentCompatibility._paragraph_to_html(para, style_name)
                if para_html:
                    html_parts.append(para_html)
            
            # Process tables
            for table in doc.tables:
                table_html = DocumentCompatibility._table_to_html(table)
                if table_html:
                    html_parts.append(table_html)
            
            # Wrap in HTML document
            return DocumentCompatibility._wrap_html(html_parts)
            
        except Exception as e:
            print(f"Method 1 failed: {e}, trying fallback...")
            return DocumentCompatibility._fallback_docx_rendering(file_path)
    
    @staticmethod
    def _paragraph_to_html(para, style_name):
        """Convert a paragraph to HTML with formatting"""
        if not para.text.strip():
            return ""
        
        # Determine heading level
        heading_level = 0
        if "Heading" in style_name:
            try:
                heading_level = int(style_name.replace("Heading", "").strip())
            except:
                heading_level = 1
        
        # Build paragraph with formatting
        text_parts = []
        for run in para.runs:
            text = run.text
            if not text:
                continue
            
            # Apply formatting
            if run.bold:
                text = f"<b>{text}</b>"
            if run.italic:
                text = f"<i>{text}</i>"
            if run.underline:
                text = f"<u>{text}</u>"
            # Add color if available
            if run.font.color and run.font.color.rgb:
                color = run.font.color.rgb
                text = f'<span style="color: #{color};">{text}</span>'
            
            text_parts.append(text)
        
        full_text = "".join(text_parts)
        
        # Apply heading or paragraph styling
        if heading_level and heading_level <= 6:
            return f"<h{heading_level}>{full_text}</h{heading_level}>"
        else:
            # Check if it's a list item
            if para.text.startswith("•") or para.text.startswith("-") or para.text.startswith("◦"):
                return f'<li style="margin-left: 20px;">{full_text}</li>'
            
            # Check for numbering (like "1.", "2.", etc.)
            import re
            if re.match(r'^\d+\.', para.text):
                return f'<li style="margin-left: 20px; list-style-type: decimal;">{full_text}</li>'
            
            # Apply alignment
            alignment = para.alignment
            align_style = ""
            if alignment:
                if alignment == 1:  # Center
                    align_style = ' style="text-align: center;"'
                elif alignment == 2:  # Right
                    align_style = ' style="text-align: right;"'
                elif alignment == 3:  # Justify
                    align_style = ' style="text-align: justify;"'
            
            return f'<p{align_style}>{full_text}</p>'
    
    @staticmethod
    def _table_to_html(table):
        """Convert a table to HTML"""
        html = ['<table style="border-collapse: collapse; width: 100%; margin: 10px 0;">']
        
        for row in table.rows:
            html.append('<tr>')
            for cell in row.cells:
                cell_text = cell.text.strip()
                # Check if it's a header row (first row often is)
                is_header = row == table.rows[0]
                tag = 'th' if is_header else 'td'
                style = 'border: 1px solid #ddd; padding: 8px; text-align: left;'
                if is_header:
                    style += ' background-color: #f2f2f2; font-weight: bold;'
                html.append(f'<{tag} style="{style}">{cell_text}</{tag}>')
            html.append('</tr>')
        
        html.append('</table>')
        return "\n".join(html)
    
    @staticmethod
    def _wrap_html(content_parts):
        """Wrap content in HTML document"""
        return f"""
        <html>
        <head>
            <style>
                body {{ 
                    font-family: 'Segoe UI', Arial, sans-serif;
                    font-size: 14px;
                    line-height: 1.6;
                    padding: 20px;
                    max-width: 900px;
                    margin: 0 auto;
                    background-color: #ffffff;
                }}
                h1 {{ font-size: 24px; color: #2c3e50; border-bottom: 2px solid #3498db; padding-bottom: 10px; }}
                h2 {{ font-size: 20px; color: #2c3e50; margin-top: 20px; }}
                h3 {{ font-size: 18px; color: #2c3e50; }}
                h4 {{ font-size: 16px; color: #2c3e50; }}
                h5 {{ font-size: 14px; color: #2c3e50; }}
                h6 {{ font-size: 13px; color: #2c3e50; }}
                p {{ margin: 8px 0; }}
                ul, ol {{ padding-left: 30px; }}
                li {{ margin: 4px 0; }}
                table {{ width: 100%; border-collapse: collapse; margin: 15px 0; }}
                th {{ background-color: #f2f2f2; font-weight: bold; }}
                td, th {{ border: 1px solid #ddd; padding: 8px; text-align: left; }}
                tr:nth-child(even) {{ background-color: #f9f9f9; }}
                .highlight {{ background-color: #fff3cd; }}
                .note {{ background-color: #d1ecf1; padding: 10px; border-radius: 4px; }}
            </style>
        </head>
        <body>
            {''.join(content_parts)}
        </body>
        </html>
        """
    
    @staticmethod
    def _fallback_docx_rendering(file_path):
        """Fallback method for DOCX rendering"""
        try:
            # Try XML extraction with formatting hints
            import zipfile
            import xml.etree.ElementTree as ET
            
            with zipfile.ZipFile(file_path, 'r') as docx_zip:
                document_xml = docx_zip.read('word/document.xml')
                root = ET.fromstring(document_xml)
                
                namespaces = {
                    'w': 'http://schemas.openxmlformats.org/wordprocessingml/2006/main',
                    'w14': 'http://schemas.microsoft.com/office/word/2010/wordml',
                }
                
                # Extract paragraphs with basic formatting
                paragraphs = []
                for para in root.findall('.//w:p', namespaces):
                    texts = []
                    for text_elem in para.findall('.//w:t', namespaces):
                        if text_elem.text:
                            texts.append(text_elem.text)
                    
                    if texts:
                        paragraphs.append(" ".join(texts))
                
                # Convert to simple HTML
                html_parts = []
                for p in paragraphs:
                    # Simple formatting detection
                    if p.strip().startswith(("•", "-", "◦")):
                        html_parts.append(f'<li style="margin-left: 20px;">{p}</li>')
                    elif p.strip().startswith("1.") or p.strip().startswith("2."):
                        html_parts.append(f'<li style="margin-left: 20px; list-style-type: decimal;">{p}</li>')
                    elif len(p) > 0 and p.upper() == p and len(p) < 100:
                        # Looks like a heading
                        html_parts.append(f'<h3>{p}</h3>')
                    else:
                        html_parts.append(f'<p>{p}</p>')
                
                return DocumentCompatibility._wrap_html(html_parts)
                
        except Exception as e:
            return f"<p>Error reading DOCX: {str(e)}</p>"
    
    @staticmethod
    def render_docx_compatibility(file_path):
        """
        Main DOCX rendering method - returns HTML formatted content
        """
        return DocumentCompatibility.render_docx_with_formatting(file_path)
    
    @staticmethod
    def render_pptx_compatibility(file_path):
        """Enhanced PPTX rendering with slide structure preservation"""
        try:
            from pptx import Presentation
            prs = Presentation(file_path)
            slides_text = []
            
            for slide_num, slide in enumerate(prs.slides, 1):
                slide_content = [f"--- Slide {slide_num} ---"]
                
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        if shape.text.strip():
                            slide_content.append(shape.text)
                    
                    if hasattr(shape, "table"):
                        for row in shape.table.rows:
                            row_text = []
                            for cell in row.cells:
                                if cell.text.strip():
                                    row_text.append(cell.text.strip())
                            if row_text:
                                slide_content.append(" | ".join(row_text))
                
                slides_text.append("\n".join(slide_content))
            
            return "\n\n".join(slides_text)
            
        except Exception as e:
            return f"Error reading PPTX: {str(e)}"
    
    @staticmethod
    def render_excel_compatibility(file_path):
        """Enhanced Excel rendering with multiple engines"""
        try:
            import openpyxl
            wb = openpyxl.load_workbook(file_path, data_only=True)
            all_data = []
            
            for sheet_name in wb.sheetnames:
                sheet = wb[sheet_name]
                sheet_data = [f"--- Sheet: {sheet_name} ---"]
                
                for row in sheet.iter_rows(values_only=True):
                    row_data = [str(cell) if cell is not None else "" for cell in row]
                    if any(row_data):
                        sheet_data.append("\t".join(row_data))
                
                all_data.append("\n".join(sheet_data))
            
            return "\n\n".join(all_data)
            
        except Exception as e:
            try:
                import pandas as pd
                df = pd.read_excel(file_path, sheet_name=None)
                all_data = []
                for sheet_name, data in df.items():
                    sheet_data = [f"--- Sheet: {sheet_name} ---"]
                    sheet_data.append(data.to_string())
                    all_data.append("\n".join(sheet_data))
                return "\n\n".join(all_data)
                
            except Exception as e2:
                return f"Error reading Excel: {str(e2)}"