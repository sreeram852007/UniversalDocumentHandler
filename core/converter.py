"""
File Conversion Engine - Convert between various document formats
Supports: PDF, DOCX, XLSX, CSV, PPTX, Images, Text, HTML, ODT, EPUB, RTF
"""

import os
import subprocess
import tempfile
from pathlib import Path
from PyQt6.QtWidgets import QMessageBox
import fitz  # PyMuPDF
from docx import Document
from docx.shared import Inches, Pt
import pandas as pd
from PIL import Image
import io
import zipfile
import xml.etree.ElementTree as ET


class FileConverter:
    """Universal file converter with multiple format support"""
    
    SUPPORTED_FORMATS = {
        'pdf': {
            'extensions': ['.pdf'],
            'name': 'PDF Document',
            'icon': '📄'
        },
        'word': {
            'extensions': ['.docx', '.doc'],
            'name': 'Word Document',
            'icon': '📝'
        },
        'excel': {
            'extensions': ['.xlsx', '.xls'],
            'name': 'Excel Spreadsheet',
            'icon': '📊'
        },
        'csv': {
            'extensions': ['.csv'],
            'name': 'CSV File',
            'icon': '📈'
        },
        'powerpoint': {
            'extensions': ['.pptx', '.ppt'],
            'name': 'PowerPoint Presentation',
            'icon': '📽️'
        },
        'image': {
            'extensions': ['.jpg', '.jpeg', '.png', '.gif', '.bmp', '.webp', '.tiff', '.svg'],
            'name': 'Image',
            'icon': '🖼️'
        },
        'text': {
            'extensions': ['.txt'],
            'name': 'Text File',
            'icon': '📄'
        },
        'html': {
            'extensions': ['.html', '.htm'],
            'name': 'HTML Document',
            'icon': '🌐'
        },
        'markdown': {
            'extensions': ['.md', '.markdown'],
            'name': 'Markdown Document',
            'icon': '📝'
        },
        'epub': {
            'extensions': ['.epub'],
            'name': 'EPUB E-Book',
            'icon': '📚'
        },
        'rtf': {
            'extensions': ['.rtf'],
            'name': 'RTF Document',
            'icon': '📄'
        },
        'odt': {
            'extensions': ['.odt'],
            'name': 'ODT Document',
            'icon': '📄'
        },
        'ods': {
            'extensions': ['.ods'],
            'name': 'ODS Spreadsheet',
            'icon': '📊'
        },
    }
    
    @staticmethod
    def get_available_conversions(input_file, file_type):
        """Get available conversion formats for a given file type"""
        conversion_map = {
            'pdf': ['word', 'text', 'html', 'epub', 'image'],
            'word': ['pdf', 'text', 'html', 'epub', 'rtf', 'odt'],
            'excel': ['csv', 'pdf', 'html', 'ods'],
            'csv': ['excel', 'pdf', 'html'],
            'powerpoint': ['pdf', 'image'],
            'image': ['pdf', 'text', 'html'],
            'text': ['pdf', 'word', 'html'],
            'html': ['pdf', 'word', 'text'],
            'markdown': ['pdf', 'html', 'text'],
            'epub': ['pdf', 'text', 'html'],
            'rtf': ['pdf', 'word', 'text'],
            'odt': ['pdf', 'word', 'text'],
            'ods': ['excel', 'csv', 'pdf'],
        }
        
        return conversion_map.get(file_type, [])
    
    @staticmethod
    def get_extension_for_format(format_name):
        """Get the default extension for a format"""
        format_extensions = {
            'pdf': '.pdf',
            'word': '.docx',
            'excel': '.xlsx',
            'csv': '.csv',
            'powerpoint': '.pptx',
            'image': '.png',
            'text': '.txt',
            'html': '.html',
            'markdown': '.md',
            'epub': '.epub',
            'rtf': '.rtf',
            'odt': '.odt',
            'ods': '.ods',
        }
        return format_extensions.get(format_name, '.txt')
    
    @staticmethod
    def convert_file(input_path, output_path):
        """Main conversion dispatcher"""
        input_path = Path(input_path)
        output_path = Path(output_path)
        
        input_ext = input_path.suffix.lower()
        output_ext = output_path.suffix.lower()
        
        # Determine file types
        input_type = None
        output_type = None
        
        for fmt, info in FileConverter.SUPPORTED_FORMATS.items():
            if input_ext in info['extensions']:
                input_type = fmt
            if output_ext in info['extensions']:
                output_type = fmt
        
        if not input_type or not output_type:
            raise ValueError(f"Unsupported conversion: {input_ext} to {output_ext}")
        
        # Dispatch to specific converter
        converter_method = f"convert_{input_type}_to_{output_type}"
        
        if hasattr(FileConverter, converter_method):
            getattr(FileConverter, converter_method)(str(input_path), str(output_path))
        else:
            # Try generic conversion methods
            if input_type == 'pdf':
                FileConverter._convert_pdf_generic(str(input_path), str(output_path), output_type)
            elif output_type == 'pdf':
                FileConverter._convert_to_pdf(str(input_path), str(output_path), input_type)
            else:
                # Try text-based conversion
                FileConverter._convert_via_text(str(input_path), str(output_path), input_type, output_type)
    
    # ============ PDF to Various Formats ============
    
    @staticmethod
    def convert_pdf_to_word(input_path, output_path):
        """Convert PDF to Word (DOCX)"""
        try:
            from pdf2docx import Converter
            cv = Converter(input_path)
            cv.convert(output_path, start=0, end=None)
            cv.close()
        except ImportError:
            # Fallback: Extract text to DOCX
            doc = fitz.open(input_path)
            word_doc = Document()
            for page in doc:
                text = page.get_text()
                if text.strip():
                    word_doc.add_paragraph(text)
            word_doc.save(output_path)
            doc.close()
    
    @staticmethod
    def convert_pdf_to_text(input_path, output_path):
        """Convert PDF to Text"""
        doc = fitz.open(input_path)
        text = ""
        for page in doc:
            text += page.get_text() + "\n\n"
        doc.close()
        
        with open(output_path, 'w', encoding='utf-8') as f:
            f.write(text)
    
    @staticmethod
    def convert_pdf_to_html(input_path, output_path):
        """Convert PDF to HTML"""
        doc = fitz.open(input_path)
        html = "<html><head><title>PDF Document</title></head><body>\n"
        for i, page in enumerate(doc):
            text = page.get_text()
            html += f"<h1>Page {i+1}</h1>\n"
            html += f"<p>{text.replace(chr(10), '<br>')}</p>\n"
        html += "</body></html>"
        doc.close()
        
        with open(output_path, 'w', encoding='utf-8') as f:
            f.write(html)
    
    @staticmethod
    def convert_pdf_to_image(input_path, output_path):
        """Convert PDF to Image (PNG)"""
        doc = fitz.open(input_path)
        page = doc[0]
        pix = page.get_pixmap(matrix=fitz.Matrix(2, 2))
        pix.save(output_path)
        doc.close()
    
    @staticmethod
    def convert_pdf_to_epub(input_path, output_path):
        """Convert PDF to EPUB"""
        # Simple EPUB generation from PDF text
        doc = fitz.open(input_path)
        text = ""
        for page in doc:
            text += page.get_text() + "\n\n"
        doc.close()
        
        # Create basic EPUB structure
        import zipfile
        import tempfile
        import os
        from datetime import datetime
        
        with tempfile.TemporaryDirectory() as tmpdir:
            # Create mimetype
            with open(os.path.join(tmpdir, 'mimetype'), 'w') as f:
                f.write('application/epub+zip')
            
            # Create META-INF/container.xml
            os.makedirs(os.path.join(tmpdir, 'META-INF'))
            with open(os.path.join(tmpdir, 'META-INF', 'container.xml'), 'w') as f:
                f.write('''<?xml version="1.0" encoding="UTF-8"?>
<container version="1.0" xmlns="urn:oasis:names:tc:opendocument:xmlns:container">
  <rootfiles>
    <rootfile full-path="book.opf" media-type="application/oebps-package+xml"/>
  </rootfiles>
</container>''')
            
            # Create book.opf
            with open(os.path.join(tmpdir, 'book.opf'), 'w') as f:
                f.write(f'''<?xml version="1.0" encoding="UTF-8"?>
<package xmlns="http://www.idpf.org/2007/opf" version="2.0" unique-identifier="BookId">
  <metadata xmlns:dc="http://purl.org/dc/elements/1.1/">
    <dc:title>PDF Document</dc:title>
    <dc:language>en</dc:language>
    <dc:date>{datetime.now().strftime('%Y-%m-%d')}</dc:date>
  </metadata>
  <manifest>
    <item id="content" href="content.html" media-type="application/xhtml+xml"/>
    <item id="nav" href="toc.ncx" media-type="application/x-dtbncx+xml"/>
  </manifest>
  <spine>
    <itemref idref="content"/>
  </spine>
</package>''')
            
            # Create content.html
            html_content = f'''<?xml version="1.0" encoding="UTF-8"?>
<!DOCTYPE html>
<html xmlns="http://www.w3.org/1999/xhtml">
<head>
  <title>PDF Document</title>
</head>
<body>
  <h1>PDF Document</h1>
  <pre>{text[:50000]}</pre>
  <p>[...truncated]</p>
</body>
</html>'''
            with open(os.path.join(tmpdir, 'content.html'), 'w', encoding='utf-8') as f:
                f.write(html_content)
            
            # Create toc.ncx
            with open(os.path.join(tmpdir, 'toc.ncx'), 'w') as f:
                f.write('''<?xml version="1.0" encoding="UTF-8"?>
<ncx xmlns="http://www.daisy.org/z3986/2005/ncx/" version="2005-1">
  <head>
    <meta name="dtb:uid" content="bookid"/>
    <meta name="dtb:depth" content="1"/>
    <meta name="dtb:generator" content="Converter"/>
  </head>
  <docTitle><text>PDF Document</text></docTitle>
  <navMap>
    <navPoint id="nav1" playOrder="1">
      <navLabel><text>Document</text></navLabel>
      <content src="content.html"/>
    </navPoint>
  </navMap>
</ncx>''')
            
            # Create EPUB zip
            with zipfile.ZipFile(output_path, 'w', zipfile.ZIP_DEFLATED) as epub:
                epub.write(os.path.join(tmpdir, 'mimetype'), 'mimetype', compress_type=zipfile.ZIP_STORED)
                epub.write(os.path.join(tmpdir, 'META-INF', 'container.xml'), 'META-INF/container.xml')
                epub.write(os.path.join(tmpdir, 'book.opf'), 'book.opf')
                epub.write(os.path.join(tmpdir, 'content.html'), 'content.html')
                epub.write(os.path.join(tmpdir, 'toc.ncx'), 'toc.ncx')
    
    # ============ Word to Various Formats ============
    
    @staticmethod
    def convert_word_to_pdf(input_path, output_path):
        """Convert Word to PDF"""
        doc = Document(input_path)
        # Simple PDF via PyMuPDF
        import tempfile
        with tempfile.NamedTemporaryFile(suffix='.txt', delete=False, mode='w', encoding='utf-8') as f:
            for para in doc.paragraphs:
                if para.text.strip():
                    f.write(para.text + '\n')
            temp_file = f.name
        
        try:
            # Create PDF from text
            doc_pdf = fitz.open()
            page = doc_pdf.new_page()
            page.insert_text((50, 50), open(temp_file, 'r', encoding='utf-8').read())
            doc_pdf.save(output_path)
            doc_pdf.close()
        finally:
            os.unlink(temp_file)
    
    @staticmethod
    def convert_word_to_text(input_path, output_path):
        """Convert Word to Text"""
        doc = Document(input_path)
        text = "\n".join([p.text for p in doc.paragraphs if p.text.strip()])
        with open(output_path, 'w', encoding='utf-8') as f:
            f.write(text)
    
    @staticmethod
    def convert_word_to_html(input_path, output_path):
        """Convert Word to HTML"""
        doc = Document(input_path)
        html = "<html><head><title>Word Document</title></head><body>\n"
        for para in doc.paragraphs:
            if para.text.strip():
                html += f"<p>{para.text}</p>\n"
        html += "</body></html>"
        
        with open(output_path, 'w', encoding='utf-8') as f:
            f.write(html)
    
    # ============ Excel/CSV to Various Formats ============
    
    @staticmethod
    def convert_excel_to_csv(input_path, output_path):
        """Convert Excel to CSV"""
        df = pd.read_excel(input_path)
        df.to_csv(output_path, index=False)
    
    @staticmethod
    def convert_excel_to_pdf(input_path, output_path):
        """Convert Excel to PDF"""
        df = pd.read_excel(input_path)
        # Create PDF with table
        doc = fitz.open()
        page = doc.new_page()
        
        # Convert dataframe to text table
        text = df.to_string()
        page.insert_text((50, 50), text[:5000])  # Limit text for PDF
        
        doc.save(output_path)
        doc.close()
    
    @staticmethod
    def convert_csv_to_excel(input_path, output_path):
        """Convert CSV to Excel"""
        df = pd.read_csv(input_path)
        df.to_excel(output_path, index=False)
    
    @staticmethod
    def convert_csv_to_pdf(input_path, output_path):
        """Convert CSV to PDF"""
        df = pd.read_csv(input_path)
        doc = fitz.open()
        page = doc.new_page()
        text = df.to_string()
        page.insert_text((50, 50), text[:5000])
        doc.save(output_path)
        doc.close()
    
    # ============ Image to Various Formats ============
    
    @staticmethod
    def convert_image_to_pdf(input_path, output_path):
        """Convert Image to PDF"""
        img = Image.open(input_path)
        # Convert to RGB if necessary
        if img.mode in ('RGBA', 'LA', 'P'):
            img = img.convert('RGB')
        img.save(output_path, 'PDF', resolution=100.0)
    
    @staticmethod
    def convert_image_to_text(input_path, output_path):
        """Convert Image to Text (Simple placeholder - would need OCR)"""
        with open(output_path, 'w', encoding='utf-8') as f:
            f.write(f"[Image to Text Conversion]\n\n")
            f.write(f"File: {os.path.basename(input_path)}\n")
            f.write(f"Size: {os.path.getsize(input_path)} bytes\n")
            f.write(f"\nNote: OCR functionality would extract text from this image.")
            f.write(f"\nFor proper OCR, consider using Tesseract OCR engine.")
    
    # ============ Generic Converters ============
    
    @staticmethod
    def _convert_pdf_generic(input_path, output_path, output_type):
        """Generic conversion from PDF"""
        converters = {
            'word': FileConverter.convert_pdf_to_word,
            'text': FileConverter.convert_pdf_to_text,
            'html': FileConverter.convert_pdf_to_html,
            'epub': FileConverter.convert_pdf_to_epub,
            'image': FileConverter.convert_pdf_to_image,
        }
        if output_type in converters:
            converters[output_type](input_path, output_path)
        else:
            raise ValueError(f"Cannot convert PDF to {output_type}")
    
    @staticmethod
    def _convert_to_pdf(input_path, output_path, input_type):
        """Generic conversion to PDF"""
        converters = {
            'word': FileConverter.convert_word_to_pdf,
            'excel': FileConverter.convert_excel_to_pdf,
            'csv': FileConverter.convert_csv_to_pdf,
            'image': FileConverter.convert_image_to_pdf,
            'text': FileConverter._convert_text_to_pdf,
            'html': FileConverter._convert_html_to_pdf,
        }
        if input_type in converters:
            converters[input_type](input_path, output_path)
        else:
            raise ValueError(f"Cannot convert {input_type} to PDF")
    
    @staticmethod
    def _convert_text_to_pdf(input_path, output_path):
        """Convert Text to PDF"""
        doc = fitz.open()
        page = doc.new_page()
        with open(input_path, 'r', encoding='utf-8', errors='ignore') as f:
            text = f.read()
        page.insert_text((50, 50), text[:10000])
        doc.save(output_path)
        doc.close()
    
    @staticmethod
    def _convert_html_to_pdf(input_path, output_path):
        """Convert HTML to PDF (simple)"""
        import re
        doc = fitz.open()
        page = doc.new_page()
        with open(input_path, 'r', encoding='utf-8', errors='ignore') as f:
            html = f.read()
        # Simple HTML tag stripping
        text = re.sub(r'<[^>]+>', '', html)
        page.insert_text((50, 50), text[:10000])
        doc.save(output_path)
        doc.close()
    
    @staticmethod
    def _convert_via_text(input_path, output_path, input_type, output_type):
        """Fallback conversion via text extraction"""
        try:
            # Extract text based on input type
            text = ""
            if input_type == 'word':
                doc = Document(input_path)
                text = "\n".join([p.text for p in doc.paragraphs if p.text.strip()])
            elif input_type in ['excel', 'csv']:
                if input_type == 'excel':
                    df = pd.read_excel(input_path)
                else:
                    df = pd.read_csv(input_path)
                text = df.to_string()
            elif input_type == 'powerpoint':
                from pptx import Presentation
                prs = Presentation(input_path)
                for slide in prs.slides:
                    for shape in slide.shapes:
                        if hasattr(shape, "text") and shape.text.strip():
                            text += shape.text + "\n"
            else:
                with open(input_path, 'r', encoding='utf-8', errors='ignore') as f:
                    text = f.read()
            
            # Write to output based on output type
            if output_type == 'text':
                with open(output_path, 'w', encoding='utf-8') as f:
                    f.write(text)
            elif output_type == 'word':
                doc = Document()
                for line in text.split('\n'):
                    if line.strip():
                        doc.add_paragraph(line.strip())
                doc.save(output_path)
            elif output_type == 'pdf':
                doc = fitz.open()
                page = doc.new_page()
                page.insert_text((50, 50), text[:10000])
                doc.save(output_path)
                doc.close()
            else:
                raise ValueError(f"Unsupported via-text conversion to {output_type}")
                
        except Exception as e:
            raise ValueError(f"Conversion failed: {str(e)}")